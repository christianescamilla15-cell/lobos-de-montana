from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# --------------- COLORS ---------------
BLACK6 = RGBColor(0x16, 0x1F, 0x28)
BLUE655 = RGBColor(0x1E, 0x37, 0x65)
BLUE2717 = RGBColor(0xA9, 0xC0, 0xDE)
GRAY642 = RGBColor(0xD3, 0xDF, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GOLD = RGBColor(0xD4, 0xA8, 0x4B)
DARK_OVERLAY = RGBColor(0x0D, 0x14, 0x1C)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(SCRIPT_DIR, "logo-wolf.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --------------- HELPERS ---------------
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri", italic=False,
             anchor=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
        # vertical centering
        txBody = txBox._element.txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', anchor)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, default_size=14,
                       default_color=WHITE):
    """Add a textbox with multiple paragraphs (list of dicts)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_info.get("text", "")
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.font.italic = line_info.get("italic", False)
        p.font.name = line_info.get("font", "Calibri")
        p.alignment = line_info.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space_after", 4))
    return txBox

def _set_shape_alpha(shape, alpha):
    """Set fill transparency. alpha: 0=fully transparent, 100=fully opaque."""
    from lxml import etree
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.find(qn('xdr:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    alpha_elem = srgb.find(qn('a:alpha'))
    if alpha_elem is None:
        alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
    alpha_elem.set('val', str(int(alpha * 1000)))

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape

def add_rounded_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape

def add_circle(slide, left, top, size, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(left), Inches(top),
                                   Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape

def add_line_shape(slide, left, top, width, color=BLUE2717, thickness=0.04):
    add_rect(slide, left, top, width, thickness, color)

def add_gradient_overlay(slide):
    """Create a gradient-like effect with layered semi-transparent rectangles."""
    # Bottom dark band
    add_rect(slide, 0, 6.0, 13.333, 1.5, DARK_OVERLAY, alpha=40)
    # Top subtle band
    add_rect(slide, 0, 0, 13.333, 0.6, BLUE655, alpha=20)
    # Right accent strip
    add_rect(slide, 12.8, 0, 0.533, 7.5, BLUE2717, alpha=10)

def footer_bar(slide):
    add_rect(slide, 0, 7.0, 13.333, 0.5, BLUE655)
    add_text(slide, 0.5, 7.05, 12.3, 0.4,
             "Lobos de Montana  |  lobosdemontana.mx  |  Conquista cada sendero",
             size=9, color=BLUE2717, align=PP_ALIGN.CENTER)

def side_bar(slide):
    add_rect(slide, 0, 0, 0.15, 7.5, BLUE2717)
    # Add decorative dots along sidebar
    for dot_y in [1.0, 2.5, 4.0, 5.5]:
        add_circle(slide, 0.03, dot_y, 0.09, WHITE, alpha=30)

def add_decorative_dots(slide, x, y, cols=5, rows=3, spacing=0.18, color=BLUE2717, alpha=15):
    """Add a grid of small decorative dots."""
    for r in range(rows):
        for c in range(cols):
            add_circle(slide, x + c * spacing, y + r * spacing, 0.06, color, alpha=alpha)

def add_logo(slide, left, top, width, height=None):
    """Add logo image to slide."""
    if not os.path.exists(LOGO_PATH):
        print(f"  [Warning] Logo not found at {LOGO_PATH}, skipping logo placement.")
        return None
    if height:
        pic = slide.shapes.add_picture(LOGO_PATH, Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    else:
        pic = slide.shapes.add_picture(LOGO_PATH, Inches(left), Inches(top),
                                       Inches(width))
    return pic

def add_logo_small(slide):
    """Add small logo in bottom-right corner."""
    add_logo(slide, 12.1, 6.2, 0.7, 0.7)

def add_section_header(slide, text):
    """Adds a styled section header with accent line."""
    add_rect(slide, 0.5, 0.3, 0.06, 0.5, BLUE2717)
    add_text(slide, 0.8, 0.3, 6, 0.6, text, size=16, color=BLUE2717, bold=True,
             font_name="Calibri Light")

def add_card_border_top(slide, left, top, width, color):
    """Adds a colored top-border accent to a card."""
    add_rect(slide, left, top, width, 0.06, color)


# ============================================================
#                    SLIDE 1 - TITLE / INTRO
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, BLACK6)

# Gradient overlays for depth
add_rect(s1, 0, 0, 13.333, 3.5, BLUE655, alpha=12)
add_rect(s1, 0, 5.5, 13.333, 2.0, BLUE655, alpha=18)

# Decorative elements
add_decorative_dots(s1, 1.0, 0.5, cols=6, rows=4, color=BLUE2717, alpha=10)
add_decorative_dots(s1, 10.5, 5.0, cols=6, rows=4, color=BLUE2717, alpha=10)

# Large decorative circles in corners
add_circle(s1, -0.5, -0.5, 2.0, BLUE655, alpha=8)
add_circle(s1, 11.8, 5.8, 2.5, BLUE655, alpha=8)

# Logo centered above title
add_logo(s1, 5.4, 0.3, 2.5, 2.5)

# Accent lines
add_line_shape(s1, 3.5, 3.0, 6.3, BLUE2717)
add_rect(s1, 5.5, 3.08, 2.3, 0.03, ACCENT_GOLD)

# Title
add_text(s1, 0.5, 3.3, 12.3, 1.5, "LOBOS DE MONTANA", size=68, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Bottom accent line
add_line_shape(s1, 3.5, 4.9, 6.3, BLUE2717)
add_rect(s1, 5.5, 4.98, 2.3, 0.03, ACCENT_GOLD)

# Tagline
add_text(s1, 1.5, 5.2, 10.3, 0.8,
         "Conquista cada sendero \u2014 equipo que resiste como tu",
         size=24, color=BLUE2717, align=PP_ALIGN.CENTER, italic=True)

# Subtitle
add_text(s1, 1.5, 6.0, 10.3, 0.6,
         "Equipo Premium de Senderismo y Montanismo",
         size=15, color=GRAY642, align=PP_ALIGN.CENTER)

footer_bar(s1)


# ============================================================
#                  SLIDE 2 - NUESTRO SLOGAN
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, BLACK6)
add_gradient_overlay(s2)
side_bar(s2)

add_section_header(s2, "NUESTRO SLOGAN")

# Main quote
add_text(s2, 0.8, 1.1, 8, 1.8,
         "\"Conquista cada sendero \u2014\nequipo que resiste como tu\"",
         size=38, color=WHITE, bold=True, font_name="Calibri Light", italic=True)

add_line_shape(s2, 0.8, 3.1, 4, BLUE2717)
add_rect(s2, 0.8, 3.18, 1.5, 0.03, ACCENT_GOLD)

items = [
    ("CONQUISTA",
     "No somos espectadores. Cada producto esta disenado para quienes se atreven a ir mas alla, a conquistar nuevos caminos y superar sus propios limites.",
     "\u2694"),
    ("CADA SENDERO",
     "No importa si es tu primera caminata o tu expedicion mas ambiciosa. Nuestro equipo se adapta a todos los terrenos y niveles de experiencia.",
     "\u26F0"),
    ("EQUIPO QUE RESISTE COMO TU",
     "La montana no perdona. Por eso seleccionamos productos con la misma resistencia y determinacion que caracteriza a nuestros clientes.",
     "\u26A1")
]
y = 3.5
for title, desc, icon in items:
    # Card background
    add_rounded_rect(s2, 0.6, y - 0.1, 11.8, 1.05, BLUE655, alpha=60)
    # Left accent bar
    add_rect(s2, 0.6, y - 0.1, 0.06, 1.05, BLUE2717)
    # Icon circle
    add_circle(s2, 0.85, y + 0.05, 0.55, BLUE655, alpha=80)
    add_text(s2, 0.85, y + 0.05, 0.55, 0.55, icon, size=20, color=BLUE2717,
             align=PP_ALIGN.CENTER)
    # Text
    add_text(s2, 1.6, y - 0.05, 10.5, 0.35, title, size=14, color=BLUE2717, bold=True)
    add_text(s2, 1.6, y + 0.32, 10.5, 0.55, desc, size=11, color=GRAY642)
    y += 1.15

# Decorative dots
add_decorative_dots(s2, 11.0, 1.0, cols=4, rows=6, alpha=12)

footer_bar(s2)
add_logo_small(s2)


# ============================================================
#                  SLIDE 3 - QUIENES SOMOS
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, BLACK6)
add_gradient_overlay(s3)
side_bar(s3)

add_section_header(s3, "QUIENES SOMOS")

add_text(s3, 0.8, 1.0, 11.5, 1.0,
         "Somos Lobos de Montana, una tienda especializada en equipo de senderismo y montanismo. Nacimos de la pasion por explorar los senderos mas desafiantes y la necesidad de contar con equipo confiable que resista cada aventura.",
         size=15, color=GRAY642)

add_line_shape(s3, 0.8, 2.2, 6, BLUE2717)

# Mission card
add_rounded_rect(s3, 0.6, 2.5, 5.8, 2.0, BLUE655, alpha=70)
add_card_border_top(s3, 0.6, 2.5, 5.8, BLUE2717)
add_circle(s3, 1.0, 2.7, 0.5, BLUE2717, alpha=30)
add_text(s3, 1.05, 2.72, 0.5, 0.5, "\u2605", size=18, color=BLUE2717, align=PP_ALIGN.CENTER)
add_text(s3, 1.7, 2.7, 4.5, 0.4, "MISION", size=16, color=BLUE2717, bold=True)
add_text(s3, 1.0, 3.3, 5.0, 1.0,
         "Equipar a cada aventurero con productos de alta calidad que garanticen seguridad, comodidad y rendimiento en cualquier terreno.",
         size=12, color=GRAY642)

# Vision card
add_rounded_rect(s3, 6.9, 2.5, 5.8, 2.0, BLUE655, alpha=70)
add_card_border_top(s3, 6.9, 2.5, 5.8, BLUE2717)
add_circle(s3, 7.3, 2.7, 0.5, BLUE2717, alpha=30)
add_text(s3, 7.35, 2.72, 0.5, 0.5, "\u25B2", size=18, color=BLUE2717, align=PP_ALIGN.CENTER)
add_text(s3, 8.0, 2.7, 4.5, 0.4, "VISION", size=16, color=BLUE2717, bold=True)
add_text(s3, 7.3, 3.3, 5.0, 1.0,
         "Ser la tienda lider en Mexico de equipo outdoor, reconocida por nuestra calidad, innovacion y compromiso con la comunidad aventurera.",
         size=12, color=GRAY642)

# Values section
add_text(s3, 0.8, 4.8, 5, 0.5, "NUESTROS VALORES", size=15, color=BLUE2717, bold=True)
add_line_shape(s3, 0.8, 5.25, 3, BLUE2717)

valores = [
    ("Calidad sin\ncompromiso", "\u2714", BLUE2717),
    ("Pasion por\nla aventura", "\u2764", RGBColor(0xC7, 0x6B, 0x6B)),
    ("Resistencia y\ndurabilidad", "\u26A1", ACCENT_GOLD),
    ("Comunidad\ny confianza", "\u2605", RGBColor(0x7B, 0xB3, 0x7B))
]
xv = 0.6
for val_text, icon, icon_color in valores:
    # Card
    add_rounded_rect(s3, xv, 5.4, 2.9, 1.2, BLUE655, alpha=65)
    # Circle icon
    add_circle(s3, xv + 0.95, 5.5, 0.55, BLACK6, alpha=50)
    add_circle(s3, xv + 1.0, 5.55, 0.45, BLUE655)
    add_text(s3, xv + 1.0, 5.55, 0.45, 0.45, icon, size=16, color=icon_color,
             align=PP_ALIGN.CENTER)
    # Label
    add_text(s3, xv + 0.2, 6.15, 2.5, 0.4, val_text, size=10, color=GRAY642,
             align=PP_ALIGN.CENTER)
    xv += 3.1

footer_bar(s3)
add_logo_small(s3)


# ============================================================
#                SLIDE 4 - NUESTROS PRODUCTOS
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, BLACK6)
add_gradient_overlay(s4)
side_bar(s4)

add_section_header(s4, "NUESTROS PRODUCTOS")
add_text(s4, 0.8, 1.0, 11, 0.5,
         "Equipo seleccionado por expertos para cada tipo de aventura",
         size=15, color=GRAY642)

# Productos Destacados label
add_text(s4, 0.8, 1.6, 5, 0.4, "PRODUCTOS DESTACADOS", size=13, color=BLUE2717, bold=True)
add_line_shape(s4, 0.8, 2.0, 11.5, BLUE2717)

dest = [
    ("Botas de Cana Media", "$1,299 MXN",
     "Soporte y traccion para terrenos irregulares", BLUE2717),
    ("Sudadera Impermeable", "$899 MXN",
     "Proteccion contra lluvia y viento", ACCENT_GOLD),
    ("Bastones de Senderismo", "$649 MXN",
     "Ultraligeros y ajustables", RGBColor(0x7B, 0xB3, 0x7B))
]
xd = 0.6
for name, price, desc, border_color in dest:
    # Card with colored top border
    add_rounded_rect(s4, xd, 2.2, 3.9, 1.8, BLUE655, alpha=65)
    add_card_border_top(s4, xd, 2.2, 3.9, border_color)
    # Price badge
    add_rounded_rect(s4, xd + 2.3, 2.35, 1.4, 0.35, border_color, alpha=70)
    add_text(s4, xd + 2.3, 2.35, 1.4, 0.35, price, size=11, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    # Product name
    add_text(s4, xd + 0.3, 2.4, 2.2, 0.4, name, size=16, color=WHITE, bold=True)
    # Separator
    add_line_shape(s4, xd + 0.3, 3.0, 3.3, border_color, 0.02)
    # Description
    add_text(s4, xd + 0.3, 3.15, 3.3, 0.5, desc, size=11, color=GRAY642)
    # Decorative corner element
    add_circle(s4, xd + 3.3, 3.4, 0.3, border_color, alpha=15)
    xd += 4.15

# Catalogo General
add_text(s4, 0.8, 4.3, 5, 0.4, "CATALOGO GENERAL", size=13, color=BLUE2717, bold=True)
add_line_shape(s4, 0.8, 4.65, 11.5, BLUE2717)

cat = [
    ("Mochila\nHidratacion 15L", "$1,499"),
    ("Gorra\nProteccion UV", "$349"),
    ("Pantalon\nConvertible", "$749"),
    ("Guantes\nTermicos", "$499"),
    ("Cantimplora\nTermica", "$399"),
    ("Gafas\nPolarizadas", "$599")
]
xc = 0.6
for name, price in cat:
    add_rounded_rect(s4, xc, 4.85, 1.95, 1.35, BLUE655, alpha=60)
    add_rect(s4, xc, 4.85, 1.95, 0.04, BLUE2717)
    add_text(s4, xc + 0.15, 4.95, 1.65, 0.65, name, size=11, color=WHITE, bold=True)
    # Price pill
    add_rounded_rect(s4, xc + 0.25, 5.7, 1.45, 0.3, BLUE2717, alpha=40)
    add_text(s4, xc + 0.25, 5.7, 1.45, 0.3, price, size=12, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    xc += 2.1

footer_bar(s4)
add_logo_small(s4)


# ============================================================
#                SLIDE 5 - NUESTRO ALCANCE
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, BLACK6)
add_gradient_overlay(s5)
side_bar(s5)

add_section_header(s5, "NUESTRO ALCANCE")
add_text(s5, 0.8, 1.0, 11, 0.8,
         "Conectamos con aventureros en todo Mexico a traves de nuestra plataforma digital y presencia en los principales destinos de montana.",
         size=15, color=GRAY642)

# Decorative dots
add_decorative_dots(s5, 10.0, 1.5, cols=8, rows=3, alpha=10)

stats = [
    ("500+", "Productos\nen catalogo", BLUE2717),
    ("10K+", "Clientes\nsatisfechos", ACCENT_GOLD),
    ("32", "Estados\nalcanzados", RGBColor(0x7B, 0xB3, 0x7B)),
]
xs = 0.6
for num, label, accent in stats:
    # Stat card
    add_rounded_rect(s5, xs, 2.2, 3.9, 3.4, BLUE655, alpha=60)
    add_card_border_top(s5, xs, 2.2, 3.9, accent)

    # Large number circle background
    add_circle(s5, xs + 0.95, 2.6, 2.0, accent, alpha=8)

    # Large number
    add_text(s5, xs + 0.3, 2.5, 3.3, 1.5, num, size=72, color=accent,
             bold=True, align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # Separator
    add_line_shape(s5, xs + 0.5, 4.2, 2.9, accent, 0.03)

    # Label
    add_text(s5, xs + 0.3, 4.4, 3.3, 0.8, label, size=18, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # Decorative bottom dots
    add_decorative_dots(s5, xs + 1.3, 5.1, cols=3, rows=2, spacing=0.2,
                        color=accent, alpha=20)
    xs += 4.15

# Bottom info bar
add_rounded_rect(s5, 1.5, 6.0, 10.3, 0.6, BLUE655, alpha=50)
add_text(s5, 1.5, 6.05, 10.3, 0.5,
         "\u26F0  lobosdemontana.mx  |  Tiendas en CDMX, Monterrey y Guadalajara  \u26F0",
         size=13, color=GRAY642, align=PP_ALIGN.CENTER)

footer_bar(s5)
add_logo_small(s5)


# ============================================================
#                    SLIDE 6 - CIERRE
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, BLACK6)

# Gradient overlays
add_rect(s6, 0, 0, 13.333, 3.0, BLUE655, alpha=10)
add_rect(s6, 0, 5.0, 13.333, 2.5, BLUE655, alpha=15)

# Decorative circles
add_circle(s6, -1, -1, 3.0, BLUE655, alpha=6)
add_circle(s6, 11.5, 5.5, 3.0, BLUE655, alpha=6)
add_decorative_dots(s6, 0.5, 0.5, cols=5, rows=3, alpha=8)
add_decorative_dots(s6, 11.0, 6.0, cols=5, rows=3, alpha=8)

# Logo centered above title
add_logo(s6, 5.4, 0.2, 2.5, 2.5)

# Accent lines
add_line_shape(s6, 3.0, 2.9, 7.3, BLUE2717)
add_rect(s6, 5.2, 2.98, 2.9, 0.03, ACCENT_GOLD)

# Title
add_text(s6, 0.5, 3.1, 12.3, 1.2, "LOBOS DE MONTANA", size=56, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Bottom accent line
add_line_shape(s6, 3.0, 4.4, 7.3, BLUE2717)
add_rect(s6, 5.2, 4.48, 2.9, 0.03, ACCENT_GOLD)

# Tagline
add_text(s6, 1.5, 4.6, 10.3, 0.6,
         "Conquista cada sendero \u2014 equipo que resiste como tu",
         size=20, color=BLUE2717, align=PP_ALIGN.CENTER, italic=True)

# Contact card
add_rounded_rect(s6, 3.2, 5.3, 6.9, 1.4, BLUE655, alpha=70)
add_card_border_top(s6, 3.2, 5.3, 6.9, BLUE2717)

add_text(s6, 3.5, 5.4, 6.3, 0.35, "CONTACTO", size=13, color=BLUE2717,
         bold=True, align=PP_ALIGN.CENTER)
add_text(s6, 3.5, 5.75, 6.3, 0.3,
         "wolves@tienda.com  |  +52 55 1234 5678",
         size=13, color=GRAY642, align=PP_ALIGN.CENTER)
add_text(s6, 3.5, 6.05, 6.3, 0.3, "lobosdemontana.mx",
         size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s6, 1.5, 6.5, 10.3, 0.4,
         "Gracias por acompanarnos en esta aventura!",
         size=14, color=GRAY642, align=PP_ALIGN.CENTER, italic=True)

footer_bar(s6)


# ============================================================
#           SLIDE 7 - EQUIPO DE DESARROLLO
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, BLACK6)
add_gradient_overlay(s7)
side_bar(s7)

add_section_header(s7, "EQUIPO DE DESARROLLO")
add_text(s7, 0.8, 1.0, 11, 0.5,
         "Las personas detras de Lobos de Montana",
         size=15, color=GRAY642, italic=True)

add_line_shape(s7, 0.8, 1.55, 5, BLUE2717)

team = [
    "Alejandro Gonzalez",
    "Erick Garcia",
    "Daniel Aguilar",
    "Alejandro Ruiz",
    "Angel Arturo",
    "Francisco Ortiz",
    "Chris Hernandez",
    "Edwin Sanchez",
    "Diego"
]

# Layout: 3 columns x 3 rows
col_positions = [0.6, 4.75, 8.9]
row_start = 1.9
row_height = 1.55
card_w = 3.8
card_h = 1.3

accent_colors = [BLUE2717, ACCENT_GOLD, RGBColor(0x7B, 0xB3, 0x7B)]

for idx, name in enumerate(team):
    col = idx % 3
    row = idx // 3
    x = col_positions[col]
    y = row_start + row * row_height
    accent = accent_colors[col % 3]

    # Card background
    add_rounded_rect(s7, x, y, card_w, card_h, BLUE655, alpha=60)
    # Left accent
    add_rect(s7, x, y, 0.06, card_h, accent)

    # Avatar circle
    add_circle(s7, x + 0.25, y + 0.25, 0.8, BLACK6, alpha=50)
    add_circle(s7, x + 0.3, y + 0.3, 0.7, BLUE655)
    # Initials
    initials = "".join([p[0] for p in name.split() if p])[:2].upper()
    add_text(s7, x + 0.3, y + 0.35, 0.7, 0.6, initials, size=20, color=accent,
             bold=True, align=PP_ALIGN.CENTER)

    # Name
    add_text(s7, x + 1.2, y + 0.3, 2.4, 0.4, name, size=16, color=WHITE, bold=True)

    # Role placeholder line
    add_line_shape(s7, x + 1.2, y + 0.75, 2.0, accent, 0.02)
    add_text(s7, x + 1.2, y + 0.8, 2.4, 0.3, "Desarrollador",
             size=10, color=GRAY642, italic=True)

# Bottom decorative
add_decorative_dots(s7, 5.5, 6.3, cols=8, rows=2, spacing=0.2, color=BLUE2717, alpha=12)
add_text(s7, 0.8, 6.4, 11.5, 0.4,
         "\u2605  Un equipo comprometido con la aventura y la tecnologia  \u2605",
         size=12, color=BLUE2717, align=PP_ALIGN.CENTER, italic=True)

footer_bar(s7)
add_logo_small(s7)


# ============================================================
#                        SAVE
# ============================================================
output_path = os.path.join(SCRIPT_DIR, "Lobos_de_Montana_Presentacion.pptx")
prs.save(output_path)
print(f"Presentacion creada exitosamente: {output_path}")
print(f"  - 7 diapositivas profesionales")
print(f"  - Logo incluido: {os.path.exists(LOGO_PATH)}")
